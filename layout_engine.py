import os
import io
import cv2
import torch
import numpy as np
import fitz 
from PIL import Image, ImageDraw, ImageFont
from dotenv import load_dotenv
from pdf2image import convert_from_path
import subprocess
import json
import uuid
import sys
import math
from transformers import AutoImageProcessor, TableTransformerForObjectDetection
from inference import get_model
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import docx
import comtypes.client

# Load environment variables
load_dotenv()

class HybridLayoutEngine:
    def __init__(self):
        self.device = "cuda" if torch.cuda.is_available() else "cpu"
        print(f"Initializing Hybrid Engine on {self.device}...")

        # --- A. LOAD MICROSOFT TABLE TRANSFORMER ---
        self.table_model_name = os.getenv("TABLE_MODEL_NAME")
        self.table_processor = AutoImageProcessor.from_pretrained(self.table_model_name)
        self.table_model = TableTransformerForObjectDetection.from_pretrained(self.table_model_name)
        self.table_model.to(self.device)
        self.table_conf = float(os.getenv("TABLE_CONF_THRESHOLD"))

        # --- B. LOAD HANDWRITING MODEL ---
        self.rf_api_key = os.getenv("ROBOFLOW_API_KEY")
        hw_id = os.getenv("HANDWRITING_MODEL_ID")
        self.hw_model = get_model(model_id=hw_id, api_key=self.rf_api_key)
        self.hw_conf = float(os.getenv("HANDWRITING_MODEL_CONFIDENCE"))

    # ==========================
    #      CORE DETECTORS
    # ==========================
    def detect_tables(self, image):
        inputs = self.table_processor(images=image, return_tensors="pt").to(self.device)
        with torch.no_grad():
            outputs = self.table_model(**inputs)
        target_sizes = torch.tensor([image.size[::-1]])
        results = self.table_processor.post_process_object_detection(outputs, threshold=self.table_conf, target_sizes=target_sizes)[0]
        table_blocks = []
        for score, label, box in zip(results["scores"], results["labels"], results["boxes"]):
            if self.table_model.config.id2label[label.item()] == 'table':
                table_blocks.append({
                    "type": "table",
                    "bbox": [round(i, 2) for i in box.tolist()],
                    "conf": round(score.item(), 2),
                    "source": "microsoft_detr"
                })
        return table_blocks

    def detect_handwriting(self, image):
        image_cv = np.array(image)
        image_cv = cv2.cvtColor(image_cv, cv2.COLOR_RGB2BGR)
        results = self.hw_model.infer(image_cv, confidence=self.hw_conf)[0]
        hw_blocks = []
        for pred in results.predictions:
            x, y, w, h = pred.x, pred.y, pred.width, pred.height
            hw_blocks.append({
                "type": "handwriting",
                "bbox": [x - w/2, y - h/2, x + w/2, y + h/2],
                "conf": pred.confidence,
                "source": "roboflow_inference"
            })
        return hw_blocks
        
    def get_text_from_image_ocr(self, image):
        temp_filename = f"temp_ocr_{uuid.uuid4()}.jpg"
        image.save(temp_filename)
        try:
            result = subprocess.run([sys.executable, "ocr_engine.py", temp_filename],
                                    capture_output=True, text=True, encoding='utf-8')
            if result.returncode == 0:
                output_lines = result.stdout.strip().split('\n')
                return json.loads(output_lines[-1] if output_lines else "[]")
            return []
        except Exception: return []
        finally:
            if os.path.exists(temp_filename): os.remove(temp_filename)

    # ==========================
    #      MERGE LOGIC
    # ==========================
    def calculate_iou(self, boxA, boxB):
        xA = max(boxA[0], boxB[0])
        yA = max(boxA[1], boxB[1])
        xB = min(boxA[2], boxB[2])
        yB = min(boxA[3], boxB[3])
        interArea = max(0, xB - xA) * max(0, yB - yA)
        boxAArea = (boxA[2] - boxA[0]) * (boxA[3] - boxA[1])
        boxBArea = (boxB[2] - boxB[0]) * (boxB[3] - boxB[1])
        if min(boxAArea, boxBArea) == 0: return 0
        return interArea / min(boxAArea, boxBArea)
    
    def calculate_intersection(self, boxA, boxB):
        # Returns the area of intersection
        xA = max(boxA[0], boxB[0])
        yA = max(boxA[1], boxB[1])
        xB = min(boxA[2], boxB[2])
        yB = min(boxA[3], boxB[3])
        interArea = max(0, xB - xA) * max(0, yB - yA)
        boxBArea = (boxB[2] - boxB[0]) * (boxB[3] - boxB[1])
        if boxBArea == 0: return 0
        return interArea / boxBArea

    def merge_consecutive_blocks(self, blocks):
        if not blocks: return []
        sorted_blocks = sorted(blocks, key=lambda b: b["bbox"][1])
        merged = []
        current = sorted_blocks[0]
        for next_b in sorted_blocks[1:]:
            is_text_merge = current["type"] in ["text", "title"] and next_b["type"] in ["text", "title"]
            is_same_type = (current["type"] == next_b["type"])
            curr_bottom = current["bbox"][3]
            next_top = next_b["bbox"][1]
            
            # Merge
            if (next_top - curr_bottom < 50) and (is_text_merge or is_same_type):
                current["bbox"] = [
                    min(current["bbox"][0], next_b["bbox"][0]),
                    min(current["bbox"][1], next_b["bbox"][1]),
                    max(current["bbox"][2], next_b["bbox"][2]),
                    max(current["bbox"][3], next_b["bbox"][3])
                ]
                current["content"] += "\n" + next_b["content"]
                current["type"] = "text" if is_text_merge else current["type"]
            else:
                merged.append(current)
                current = next_b
        merged.append(current)
        return merged

    def intelligent_merge(self, base_blocks, table_blocks, hw_blocks):
        final_blocks = []
        consumed_indices = set()
        
        # 1. Process Tables
        for t_block in table_blocks:
            t_content = []
            t_bbox = t_block["bbox"]
            
            # Find all base text blocks that are inside this table
            for idx, p_block in enumerate(base_blocks):
                if self.calculate_intersection(t_bbox, p_block["bbox"]) > 0.5:
                    t_content.append(p_block["text"])
                    consumed_indices.add(idx)
            
            # Join the found text to form the table content
            full_table_text = "\n".join(t_content)
            
            final_blocks.append({
                "type": "table",
                "bbox": t_bbox,
                "action": "extract_table_logic",
                "content": full_table_text if full_table_text else "[EMPTY_TABLE]"
            })

        # 2. Process Remaining Text
        for idx, p_block in enumerate(base_blocks):
            if idx in consumed_indices:
                continue
            
            p_bbox = p_block["bbox"]
            assigned_type = p_block["type"]
            action = "crop_image" if assigned_type == "image" else "extract_text"

            final_blocks.append({
                "type": assigned_type,
                "bbox": p_bbox,
                "action": action,
                "content": p_block.get("text", "")
            })

        # 3. Isolated Handwriting
        for hw in hw_blocks:
            is_isolated = True
            for fb in final_blocks:
                if self.calculate_iou(hw["bbox"], fb["bbox"]) > 0.2:
                    is_isolated = False; break
            if is_isolated:
                final_blocks.append({
                    "type": "handwriting_region", "bbox": hw["bbox"],
                    "action": "send_to_ocr", "content": "[HANDWRITING_IMAGE]"
                })
                
        return self.merge_consecutive_blocks(final_blocks)

    def analyze_single_page(self, image, base_text_blocks, page_num, output_dir, visualize=False):
        table_blocks = self.detect_tables(image)
        hw_blocks = self.detect_handwriting(image)
        merged = self.intelligent_merge(base_text_blocks, table_blocks, hw_blocks)
        if visualize: self.visualize_page(image.copy(), merged, page_num, output_dir)
        return merged

    # ==========================
    #      FILE HANDLERS
    # ==========================
    def process_pdf(self, pdf_path, output_dir, visualize=False):
        print(f"Processing PDF: {pdf_path}")
        poppler_path = os.getenv("POPPLER_PATH")
        doc = fitz.open(pdf_path)
        page_images = convert_from_path(pdf_path, poppler_path=poppler_path)
        results = []
        for page_num, (page_obj, page_img) in enumerate(zip(doc, page_images)):
            raw_blocks = page_obj.get_text("dict")["blocks"]
            base_blocks = []
            img_w, img_h = page_img.size
            pdf_w, pdf_h = page_obj.rect.width, page_obj.rect.height
            page_area = img_w * img_h
            for b in raw_blocks:
                x0, y0, x1, y1 = b['bbox']
                scaled_box = [x0*(img_w/pdf_w), y0*(img_h/pdf_h), x1*(img_w/pdf_w), y1*(img_h/pdf_h)]
                if b["type"] == 1: 
                    if ((scaled_box[2]-scaled_box[0]) * (scaled_box[3]-scaled_box[1]) / page_area) > 0.8: continue
                b_type = "image" if b["type"] == 1 else "text"
                text_content = "[IMAGE_BINARY]" if b["type"] == 1 else ""
                if b["type"] == 0:
                    try: 
                        if b["lines"][0]["spans"][0]["size"] > 14: b_type = "title"
                    except: pass
                    for line in b.get("lines", []):
                        for span in line.get("spans", []): text_content += span["text"] + " "
                base_blocks.append({"bbox": scaled_box, "type": b_type, "text": text_content.strip()})
            
            page_res = self.analyze_single_page(page_img, base_blocks, page_num+1, output_dir, visualize)
            results.append({"page": page_num+1, "blocks": page_res})
        return results

    def process_image(self, image_path, output_dir, visualize=False):
        print(f"Processing Image: {image_path}")
        image = Image.open(image_path).convert("RGB")
        print("   Running base OCR (Subprocess)...")
        base_blocks = self.get_text_from_image_ocr(image)
        page_res = self.analyze_single_page(image, base_blocks, 1, output_dir, visualize)
        return [{"page": 1, "blocks": page_res}]

    def process_pptx_structurally(self, pptx_path, output_dir):
        print(f"Processing PPTX (Structural Mode): {pptx_path}")
        try: prs = Presentation(pptx_path)
        except Exception as e:
            print(f"Failed to load PPTX structurally: {e}"); return []
        results = []
        for i, slide in enumerate(prs.slides):
            print(f"   --- Slide {i+1} ---")
            blocks = []
            slide_w, slide_h = 1280, 720 
            reconstructed_img = Image.new('RGB', (slide_w, slide_h), 'white')
            draw = ImageDraw.Draw(reconstructed_img)
            
            for shape in slide.shapes:
                if not hasattr(shape, "left"): continue
                x = int(shape.left / 9525); y = int(shape.top / 9525)
                w = int(shape.width / 9525); h = int(shape.height / 9525)
                if w <= 0 or h <= 0: continue
                bbox = [x, y, x+w, y+h]

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_blob = shape.image.blob
                        shape_img = Image.open(io.BytesIO(image_blob)).convert("RGB").resize((w, h))
                        reconstructed_img.paste(shape_img, (x, y))
                        blocks.append({"type": "image", "bbox": bbox, "action": "crop_image", "content": "[IMAGE_BINARY]"})
                    except: pass
                
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    draw.rectangle([x, y, x+w, y+h], outline="gray", width=2)
                    draw.text((x+5, y+5), "[TABLE]", fill="gray")
                    table_rows = []
                    for row in shape.table.rows:
                        row_cells = []
                        for cell in row.cells:
                            row_cells.append(cell.text_frame.text.strip())
                        table_rows.append(" | ".join(row_cells))
                    blocks.append({"type": "table", "bbox": bbox, "action": "extract_table_logic", "content": "\n".join(table_rows)})

                if hasattr(shape, "text") and shape.text.strip():
                    text_content = shape.text.strip()
                    try: font = ImageFont.truetype("arial.ttf", 14)
                    except: font = ImageFont.load_default()
                    draw.text((x, y), text_content[:50], fill="black", font=font)
                    blocks.append({"type": "text", "bbox": bbox, "action": "extract_text", "content": text_content})

            merged_blocks = self.merge_consecutive_blocks(blocks)
            if os.path.exists(output_dir): self.visualize_page(reconstructed_img, merged_blocks, i+1, output_dir)
            results.append({"page": i+1, "blocks": merged_blocks})
        return results

    def process_docx(self, docx_path, output_dir, visualize=False):
        print(f"Processing DOCX: {docx_path}")
        pdf_path = os.path.splitext(docx_path)[0] + ".pdf"
        pdf_path = os.path.abspath(pdf_path)
        docx_path = os.path.abspath(docx_path)
        converted = False
        if not os.path.exists(pdf_path):
            print("   Converting DOCX to PDF...")
            try:
                word = comtypes.client.CreateObject('Word.Application')
                doc = word.Documents.Open(docx_path)
                doc.SaveAs(pdf_path, FileFormat=17) 
                doc.Close(); word.Quit(); converted = True
            except Exception as e:
                print(f"DOCX->PDF Conversion Failed: {e}")
        else: converted = True

        if converted:
            return self.process_pdf(pdf_path, output_dir, visualize)
        else:
            print("   Switching to Structural Fallback...")
            return self.process_docx_structurally(docx_path, output_dir)

    def process_docx_structurally(self, docx_path, output_dir):
        """
        Visualizes DOCX content by drawing it onto a blank page.
        """
        try:
            doc = docx.Document(docx_path)
            blocks = []
            
            # Simple Visualization Canvas
            page_w, page_h = 800, 1000
            current_y = 50
            reconstructed_img = Image.new('RGB', (page_w, page_h), 'white')
            draw = ImageDraw.Draw(reconstructed_img)
            
            try: font = ImageFont.truetype("arial.ttf", 14)
            except: font = ImageFont.load_default()

            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    bbox = [50, current_y, 750, current_y + 20]
                    blocks.append({
                        "type": "text", "bbox": bbox, "action": "extract_text", "content": text
                    })
                    draw.text((50, current_y), text[:100], fill="black", font=font)
                    current_y += 30 
                    if current_y > 950: break

            if os.path.exists(output_dir):
                self.visualize_page(reconstructed_img, blocks, 1, output_dir)

            return [{"page": 1, "blocks": blocks}]
        except Exception as e:
            print(f"Failed to process DOCX structurally: {e}")
            return []

    def process_pptx(self, pptx_path, output_dir, visualize=False):
        pdf_path = os.path.splitext(pptx_path)[0] + ".pdf"
        pdf_path = os.path.abspath(pdf_path); pptx_path = os.path.abspath(pptx_path)
        converted = False
        if not os.path.exists(pdf_path):
            try:
                powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                deck = powerpoint.Presentations.Open(pptx_path)
                deck.SaveAs(pdf_path, 32); deck.Close(); powerpoint.Quit()
                converted = True
            except: converted = False
        else: converted = True

        if converted: return self.process_pdf(pdf_path, output_dir, visualize)
        else: return self.process_pptx_structurally(pptx_path, output_dir)

    def process_file(self, file_path, output_root="data_output", visualize=False):
        file_name = os.path.splitext(os.path.basename(file_path))[0]
        output_dir = os.path.join(output_root, file_name)
        if not os.path.exists(output_dir): os.makedirs(output_dir)
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pdf": return self.process_pdf(file_path, output_dir, visualize)
        elif ext in [".jpg", ".jpeg", ".png", ".bmp"]: return self.process_image(file_path, output_dir, visualize)
        elif ext in [".pptx", ".ppt"]: return self.process_pptx(file_path, output_dir, visualize)
        elif ext in [".docx", ".doc"]: return self.process_docx(file_path, output_dir, visualize)
        else:
            print(f"Unsupported file type: {ext}")
            return []

    def visualize_page(self, image, blocks, page_num, output_dir):
        draw = ImageDraw.Draw(image)
        try: font = ImageFont.truetype("arial.ttf", 24)
        except: font = ImageFont.load_default()
        for block in blocks:
            bbox = block["bbox"]; b_type = block["type"]
            if "handwriting" in b_type: color = "red"; width = 4
            elif "table" in b_type: color = "blue"; width = 3
            elif "title" in b_type: color = "orange"; width = 3
            elif "image" in b_type: color = "magenta"; width = 3
            else: color = "green"; width = 2
            draw.rectangle(bbox, outline=color, width=width)
        save_path = os.path.join(output_dir, f"annotated_page_{page_num}.jpg")
        image.save(save_path)
        print(f"   Saved annotation: {save_path}")